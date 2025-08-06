#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PyMuPDF Pro 统一文档处理实现示例
"""

import os
import tempfile
from pathlib import Path
import logging

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def test_pymupdf_pro_installation():
    """测试 PyMuPDF Pro 安装"""
    print("🔍 测试 PyMuPDF Pro 安装...")
    
    try:
        import pymupdf.pro
        print("✅ PyMuPDF Pro 已安装")
        
        # 尝试解锁（需要试用密钥）
        try:
            pymupdf.pro.unlock()
            print("✅ PyMuPDF Pro 解锁成功")
            return True
        except Exception as e:
            print(f"⚠️  PyMuPDF Pro 解锁失败: {e}")
            print("💡 需要获取试用密钥，请访问: https://pymupdf.cn/en/latest/pymupdf-pro.html")
            return False
            
    except ImportError:
        print("❌ PyMuPDF Pro 未安装")
        print("💡 安装命令: pip install pymupdfpro")
        return False

def create_test_documents():
    """创建测试文档"""
    print("\n📝 创建测试文档...")
    
    test_docs = {}
    
    # 创建测试 Word 文档
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('测试文档', 0)
        doc.add_paragraph('这是一个测试 Word 文档。')
        doc.add_heading('功能特点', level=1)
        doc.add_paragraph('1. 用户管理\n2. 知识管理\n3. 搜索功能')
        
        word_path = "test_document.docx"
        doc.save(word_path)
        test_docs['word'] = word_path
        print("✅ 创建测试 Word 文档")
    except Exception as e:
        print(f"❌ 创建 Word 文档失败: {e}")
    
    # 创建测试 Excel 文档
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "测试数据"
        ws['A1'] = "项目名称"
        ws['B1'] = "状态"
        ws['A2'] = "知识管理系统"
        ws['B2'] = "开发中"
        ws['A3'] = "用户管理模块"
        ws['B3'] = "已完成"
        
        excel_path = "test_spreadsheet.xlsx"
        wb.save(excel_path)
        test_docs['excel'] = excel_path
        print("✅ 创建测试 Excel 文档")
    except Exception as e:
        print(f"❌ 创建 Excel 文档失败: {e}")
    
    # 创建测试 PowerPoint 文档
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "测试演示文稿"
        subtitle.text = "这是一个测试 PowerPoint 文档"
        
        ppt_path = "test_presentation.pptx"
        prs.save(ppt_path)
        test_docs['powerpoint'] = ppt_path
        print("✅ 创建测试 PowerPoint 文档")
    except Exception as e:
        print(f"❌ 创建 PowerPoint 文档失败: {e}")
    
    return test_docs

def process_document_with_pymupdf_pro(file_path):
    """使用 PyMuPDF Pro 处理文档"""
    print(f"\n🔍 使用 PyMuPDF Pro 处理文档: {file_path}")
    
    try:
        import pymupdf.pro
        pymupdf.pro.unlock()
        
        # 打开文档
        doc = pymupdf.open(file_path)
        print(f"✅ 成功打开文档，页数: {len(doc)}")
        
        # 提取文本
        text_parts = []
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"=== 第 {page_num + 1} 页 ===\n{text}")
        
        full_text = "\n\n".join(text_parts)
        print(f"✅ 文本提取完成，总字符数: {len(full_text)}")
        
        # 显示前500个字符
        preview = full_text[:500] + "..." if len(full_text) > 500 else full_text
        print(f"📄 文本预览:\n{preview}")
        
        return full_text
        
    except Exception as e:
        print(f"❌ 文档处理失败: {e}")
        return None

def compare_processing_methods():
    """比较不同处理方法的差异"""
    print("\n📊 处理方法对比...")
    
    # 创建测试文档
    test_docs = create_test_documents()
    
    if not test_docs:
        print("❌ 无法创建测试文档，跳过对比")
        return
    
    # 对比处理结果
    for doc_type, file_path in test_docs.items():
        print(f"\n{'='*50}")
        print(f"📄 处理 {doc_type.upper()} 文档: {file_path}")
        
        # 使用 PyMuPDF Pro 处理
        pymupdf_result = process_document_with_pymupdf_pro(file_path)
        
        if pymupdf_result:
            print(f"✅ PyMuPDF Pro 处理成功")
            print(f"   文本长度: {len(pymupdf_result)} 字符")
        else:
            print(f"❌ PyMuPDF Pro 处理失败")

def demonstrate_unified_api():
    """演示统一的 API 接口"""
    print("\n🎯 演示统一的 API 接口...")
    
    try:
        import pymupdf.pro
        pymupdf.pro.unlock()
        
        # 统一的文档处理函数
        def process_any_document(file_path):
            """统一的文档处理函数"""
            doc = pymupdf.open(file_path)
            text_parts = []
            
            for page in doc:
                text = page.get_text()
                if text.strip():
                    text_parts.append(text)
            
            return "\n\n".join(text_parts)
        
        # 测试不同格式
        test_files = [
            "test_document.docx",
            "test_spreadsheet.xlsx", 
            "test_presentation.pptx"
        ]
        
        for file_path in test_files:
            if os.path.exists(file_path):
                print(f"\n📄 处理: {file_path}")
                result = process_any_document(file_path)
                print(f"   结果长度: {len(result)} 字符")
                print(f"   前100字符: {result[:100]}...")
        
        print("\n✅ 统一 API 演示完成")
        
    except Exception as e:
        print(f"❌ 统一 API 演示失败: {e}")

def show_implementation_benefits():
    """展示实现优势"""
    print("\n💡 PyMuPDF Pro 实现优势:")
    
    print("\n1. 代码简化:")
    print("   当前方案: 需要4个不同的库和4套不同的处理逻辑")
    print("   PyMuPDF Pro: 1个库，1套统一的处理逻辑")
    
    print("\n2. 维护成本:")
    print("   当前方案: 高 - 需要维护多个库的版本兼容性")
    print("   PyMuPDF Pro: 低 - 只需要维护一个库")
    
    print("\n3. 扩展性:")
    print("   当前方案: 新增格式需要大量代码修改")
    print("   PyMuPDF Pro: 新增格式只需要简单配置")
    
    print("\n4. 一致性:")
    print("   当前方案: 不同格式的文本提取质量不一致")
    print("   PyMuPDF Pro: 所有格式使用相同的提取逻辑")

def main():
    """主函数"""
    print("🚀 PyMuPDF Pro 统一文档处理演示")
    print("=" * 60)
    
    # 测试安装
    if not test_pymupdf_pro_installation():
        print("\n❌ PyMuPDF Pro 不可用，无法继续演示")
        return
    
    # 创建测试文档
    test_docs = create_test_documents()
    
    # 对比处理方法
    compare_processing_methods()
    
    # 演示统一 API
    demonstrate_unified_api()
    
    # 展示优势
    show_implementation_benefits()
    
    print("\n" + "=" * 60)
    print("📋 总结:")
    print("✅ PyMuPDF Pro 提供了统一的文档处理方案")
    print("✅ 大幅简化了代码复杂度和维护成本")
    print("✅ 建议在评估试用版后考虑采用此方案")
    
    # 清理测试文件
    print("\n🧹 清理测试文件...")
    for file_path in test_docs.values():
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"   删除: {file_path}")

if __name__ == "__main__":
    main() 