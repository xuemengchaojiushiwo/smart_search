#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
真实文档重叠优化测试脚本
使用实际的文档处理来测试重叠效果
"""

import os
import time
import json
from pathlib import Path
from typing import Dict, Any, List
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# 导入配置
import sys
sys.path.append('python_service')
from config import DOCUMENT_CONFIG, CHUNKING_CONFIG

# 尝试导入PyMuPDF Pro
try:
    import pymupdf
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

class RealOverlapOptimizationTester:
    """真实文档重叠优化测试器"""
    
    def __init__(self):
        self.test_files = [
            "python_service/file/安联美元.pdf",
            "python_service/file/店铺入住流程.pdf", 
            "python_service/file/manus介绍.docx",
            "python_service/file/小红书选品.xlsx",
            "python_service/file/network_skill.pptx"
        ]
    
    def get_dynamic_overlap(self, file_extension: str) -> int:
        """根据文件类型获取动态重叠大小"""
        extension = file_extension.lower().replace('.', '')
        overlap_ratio = DOCUMENT_CONFIG['dynamic_overlap'].get(extension, DOCUMENT_CONFIG['dynamic_overlap']['default'])
        return int(DOCUMENT_CONFIG['chunk_size'] * overlap_ratio)
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """从PDF提取文本"""
        if PYMUPDF_AVAILABLE:
            try:
                doc = pymupdf.open(file_path)
                text_parts = []
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"=== 第 {page_num + 1} 页 ===\n{text}")
                return "\n\n".join(text_parts)
            except Exception as e:
                print(f"PyMuPDF Pro 处理失败: {e}")
                return f"PDF文件内容模拟 - {os.path.basename(file_path)}"
        else:
            return f"PDF文件内容模拟 - {os.path.basename(file_path)}"
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """从DOCX提取文本"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"DOCX处理失败: {e}")
            return f"DOCX文件内容模拟 - {os.path.basename(file_path)}"
    
    def extract_text_from_xlsx(self, file_path: str) -> str:
        """从XLSX提取文本"""
        try:
            from openpyxl import load_workbook
            workbook = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"工作表: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None and str(cell_value).strip():
                            row_text.append(str(cell_value).strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"XLSX处理失败: {e}")
            return f"XLSX文件内容模拟 - {os.path.basename(file_path)}"
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """从PPTX提取文本"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                text_parts.append(f"幻灯片 {slide_number}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                text_parts.append("")  # 空行分隔幻灯片
            
            return "\n".join(text_parts)
        except Exception as e:
            print(f"PPTX处理失败: {e}")
            return f"PPTX文件内容模拟 - {os.path.basename(file_path)}"
    
    def extract_text_from_file(self, file_path: str) -> str:
        """根据文件类型提取文本"""
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif file_extension == '.docx':
            return self.extract_text_from_docx(file_path)
        elif file_extension == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif file_extension == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            return f"未知文件类型内容模拟 - {os.path.basename(file_path)}"
    
    def analyze_overlap_quality(self, chunks: List[Document]) -> Dict[str, Any]:
        """分析重叠质量"""
        analysis = {
            'total_chunks': len(chunks),
            'overlap_analysis': [],
            'coherence_score': 0,
            'issues': [],
            'strengths': []
        }
        
        if len(chunks) <= 1:
            return analysis
        
        total_overlap_score = 0
        total_issues = 0
        total_strengths = 0
        
        for i in range(len(chunks) - 1):
            chunk1 = chunks[i].page_content
            chunk2 = chunks[i + 1].page_content
            
            # 计算重叠
            overlap = self.calculate_overlap(chunk1, chunk2)
            overlap_ratio = overlap / len(chunk2) if len(chunk2) > 0 else 0
            
            # 分析重叠质量
            overlap_quality = {
                'chunk_pair': f"{i+1}-{i+2}",
                'overlap_chars': overlap,
                'overlap_ratio': overlap_ratio,
                'chunk1_end': chunk1[-100:] if len(chunk1) > 100 else chunk1,
                'chunk2_start': chunk2[:100] if len(chunk2) > 100 else chunk2,
                'quality_score': 0
            }
            
            # 评估重叠质量
            if overlap_ratio >= 0.25:  # 25%以上重叠
                overlap_quality['quality_score'] = 90
                total_strengths += 1
                analysis['strengths'].append(f"块{i+1}-{i+2}: 重叠比例{overlap_ratio:.1%}，质量良好")
            elif overlap_ratio >= 0.15:  # 15-25%重叠
                overlap_quality['quality_score'] = 70
                analysis['strengths'].append(f"块{i+1}-{i+2}: 重叠比例{overlap_ratio:.1%}，质量一般")
            else:  # 15%以下重叠
                overlap_quality['quality_score'] = 30
                total_issues += 1
                analysis['issues'].append(f"块{i+1}-{i+2}: 重叠比例{overlap_ratio:.1%}，重叠不足")
            
            total_overlap_score += overlap_quality['quality_score']
            analysis['overlap_analysis'].append(overlap_quality)
        
        # 计算总体连贯性分数
        if len(chunks) > 1:
            analysis['coherence_score'] = max(0, min(100, total_overlap_score / (len(chunks) - 1)))
        
        return analysis
    
    def calculate_overlap(self, text1: str, text2: str) -> int:
        """计算两个文本的重叠字符数"""
        words1 = text1.split()
        words2 = text2.split()
        
        if len(words1) < 3 or len(words2) < 3:
            return 0
        
        # 检查结尾和开头的重叠
        for i in range(min(20, len(words1))):  # 检查更多词
            end_words = ' '.join(words1[-(i+1):])
            if text2.startswith(end_words):
                return len(end_words)
        
        return 0
    
    def test_file_overlap(self, file_path: str) -> Dict[str, Any]:
        """测试单个文件的重叠效果"""
        print(f"测试文件: {file_path}")
        
        if not os.path.exists(file_path):
            return {'success': False, 'error': '文件不存在'}
        
        start_time = time.time()
        
        try:
            # 获取文件扩展名
            file_extension = Path(file_path).suffix.lower()
            
            # 根据文件类型获取动态重叠
            dynamic_overlap = self.get_dynamic_overlap(file_extension)
            
            # 提取文本
            text = self.extract_text_from_file(file_path)
            
            if len(text) < 100:  # 文本太短，无法有效测试
                return {
                    'success': False, 
                    'error': f'文本太短({len(text)}字符)，无法有效测试重叠效果'
                }
            
            # 使用优化后的配置进行分块
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=dynamic_overlap,  # 使用动态重叠
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
            end_time = time.time()
            
            # 分析重叠质量
            overlap_analysis = self.analyze_overlap_quality(chunks)
            
            return {
                'success': True,
                'file_path': file_path,
                'file_extension': file_extension,
                'processing_time': end_time - start_time,
                'total_chunks': len(chunks),
                'chunk_size': DOCUMENT_CONFIG['chunk_size'],
                'dynamic_overlap': dynamic_overlap,
                'overlap_ratio': dynamic_overlap / DOCUMENT_CONFIG['chunk_size'],
                'total_text_length': len(text),
                'avg_chunk_length': sum(len(chunk.page_content) for chunk in chunks) / len(chunks) if chunks else 0,
                'overlap_analysis': overlap_analysis,
                'file_size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'text_preview': text[:200] + '...' if len(text) > 200 else text
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'processing_time': time.time() - start_time
            }
    
    def run_all_tests(self) -> Dict[str, Any]:
        """运行所有测试"""
        print("开始真实文档重叠优化测试...")
        
        results = {}
        total_files = 0
        successful_files = 0
        
        for file_path in self.test_files:
            if os.path.exists(file_path):
                total_files += 1
                result = self.test_file_overlap(file_path)
                results[Path(file_path).name] = result
                
                if result['success']:
                    successful_files += 1
                    coherence_score = result['overlap_analysis']['coherence_score']
                    overlap_ratio = result['overlap_ratio']
                    total_chunks = result['total_chunks']
                    print(f"✅ {Path(file_path).name}: 重叠比例 {overlap_ratio:.1%}, 连贯性分数 {coherence_score:.1f}, 块数 {total_chunks}")
                else:
                    print(f"❌ {Path(file_path).name}: {result['error']}")
            else:
                print(f"⚠️ 文件不存在: {file_path}")
        
        # 生成总结报告
        summary = {
            'total_files': total_files,
            'successful_files': successful_files,
            'success_rate': successful_files / total_files if total_files > 0 else 0,
            'average_coherence_score': 0,
            'average_overlap_ratio': 0,
            'total_chunks_generated': 0
        }
        
        if successful_files > 0:
            coherence_scores = [r['overlap_analysis']['coherence_score'] for r in results.values() if r['success']]
            overlap_ratios = [r['overlap_ratio'] for r in results.values() if r['success']]
            total_chunks = [r['total_chunks'] for r in results.values() if r['success']]
            
            summary['average_coherence_score'] = sum(coherence_scores) / len(coherence_scores)
            summary['average_overlap_ratio'] = sum(overlap_ratios) / len(overlap_ratios)
            summary['total_chunks_generated'] = sum(total_chunks)
        
        results['summary'] = summary
        
        return results

def main():
    """主函数"""
    tester = RealOverlapOptimizationTester()
    results = tester.run_all_tests()
    
    # 保存结果
    output_file = "real_overlap_optimization_test_report.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"\n测试完成！结果已保存到: {output_file}")
    print(f"总结:")
    print(f"- 测试文件数: {results['summary']['total_files']}")
    print(f"- 成功处理: {results['summary']['successful_files']}")
    print(f"- 成功率: {results['summary']['success_rate']:.1%}")
    print(f"- 平均连贯性分数: {results['summary']['average_coherence_score']:.1f}")
    print(f"- 平均重叠比例: {results['summary']['average_overlap_ratio']:.1%}")
    print(f"- 总生成块数: {results['summary']['total_chunks_generated']}")

if __name__ == "__main__":
    main()
