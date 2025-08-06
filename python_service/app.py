from fastapi import FastAPI, HTTPException, UploadFile, File
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import uvicorn
import json
import logging
import os
import tempfile
from pathlib import Path

# 文档处理相关
from langchain.document_loaders import PyMuPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import ElasticsearchStore
from langchain.schema import Document

# 新增文档处理库
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# ES相关
from elasticsearch import Elasticsearch
import hashlib

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Knowledge Base Python Service", version="1.0.0")

# 导入配置
from config import ES_CONFIG, DOCUMENT_CONFIG, EMBEDDING_CONFIG, RAG_CONFIG

# 初始化ES客户端
es_client = Elasticsearch(
    [f"http://{ES_CONFIG['host']}:{ES_CONFIG['port']}"],
    basic_auth=(ES_CONFIG['username'], ES_CONFIG['password']) if ES_CONFIG['username'] else None,
    verify_certs=ES_CONFIG['verify_certs']
)

# 初始化embedding模型
embeddings = HuggingFaceEmbeddings(
    model_name=EMBEDDING_CONFIG['model_name'],
    model_kwargs={'device': EMBEDDING_CONFIG['device']}
)

# 请求模型
class LdapValidateRequest(BaseModel):
    username: str
    password: str

class ChatRequest(BaseModel):
    question: str
    user_id: str

class DocumentProcessRequest(BaseModel):
    knowledge_id: int
    knowledge_name: str
    description: str
    tags: List[str]
    effective_time: str

# 响应模型
class LdapValidateResponse(BaseModel):
    success: bool
    message: str
    email: Optional[str] = None
    role: Optional[str] = None

class KnowledgeReference(BaseModel):
    knowledge_id: int
    knowledge_name: str
    description: str
    tags: List[str]
    effective_time: str
    attachments: List[str]
    relevance: float

class ChatResponse(BaseModel):
    answer: str
    references: List[KnowledgeReference]
    session_id: Optional[str] = None

class DocumentProcessResponse(BaseModel):
    success: bool
    message: str
    chunks_count: int
    knowledge_id: int

@app.get("/")
def read_root():
    return {"message": "Knowledge Base Python Service"}

@app.post("/api/ldap/validate", response_model=LdapValidateResponse)
def validate_ldap_user(request: LdapValidateRequest):
    """
    LDAP用户验证
    这里应该调用实际的LDAP服务
    """
    logger.info(f"LDAP验证请求: {request.username}")
    
    try:
        # 模拟LDAP验证逻辑
        # 实际项目中应该调用真实的LDAP服务
        if request.username == "admin" and request.password == "password":
            return LdapValidateResponse(
                success=True,
                message="验证成功",
                email="admin@example.com",
                role="admin"
            )
        elif request.username == "user" and request.password == "password":
            return LdapValidateResponse(
                success=True,
                message="验证成功",
                email="user@example.com",
                role="user"
            )
        else:
            return LdapValidateResponse(
                success=False,
                message="用户名或密码错误"
            )
    except Exception as e:
        logger.error(f"LDAP验证失败: {str(e)}")
        return LdapValidateResponse(
            success=False,
            message=f"LDAP验证服务异常: {str(e)}"
        )

@app.post("/api/document/process", response_model=DocumentProcessResponse)
async def process_document(
    file: UploadFile = File(...),
    knowledge_id: int = None,
    knowledge_name: str = None,
    description: str = None,
    tags: str = None,
    effective_time: str = None
):
    """
    处理上传的文档（PDF、WORD、EXCEL、TXT）
    使用langchain+pymupdf4llm解析并存入ES
    """
    logger.info(f"开始处理文档: {file.filename}, 知识ID: {knowledge_id}")
    
    try:
        # 检查文件类型
        file_extension = Path(file.filename).suffix.lower()
        
        if file_extension not in DOCUMENT_CONFIG['allowed_extensions']:
            raise HTTPException(status_code=400, detail=f"不支持的文件类型: {file_extension}")
        
        # 保存上传的文件
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        # 解析文档
        chunks = []
        if file_extension == '.pdf':
            # 使用PyMuPDFLoader处理PDF
            loader = PyMuPDFLoader(temp_file_path)
            documents = loader.load()
            
            # 文本分割
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_documents(documents)
            
        elif file_extension == '.txt':
            # 处理TXT文件
            with open(temp_file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
        
        elif file_extension == '.docx':
            # 处理Word文档
            doc = DocxDocument(temp_file_path)
            text_parts = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
        elif file_extension == '.xlsx':
            # 处理Excel文件
            workbook = load_workbook(temp_file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"工作表: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None and str(cell_value).strip():
                            row_text.append(str(cell_value).strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
        elif file_extension == '.pptx':
            # 处理PowerPoint文件
            prs = Presentation(temp_file_path)
            text_parts = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                text_parts.append(f"幻灯片 {slide_number}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                text_parts.append("")  # 空行分隔幻灯片
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
        
        # 为每个chunk添加元数据
        for chunk in chunks:
            chunk.metadata.update({
                "knowledge_id": knowledge_id,
                "knowledge_name": knowledge_name,
                "description": description,
                "tags": tags.split(",") if tags else [],
                "effective_time": effective_time,
                "source_file": file.filename,
                "file_type": file_extension
            })
        
        # 存入ES
        vectorstore = ElasticsearchStore(
            embedding=embeddings,
            index_name=ES_CONFIG['index'],
            es_connection=es_client
        )
        
        # 添加文档到向量存储
        vectorstore.add_documents(chunks)
        
        # 清理临时文件
        os.unlink(temp_file_path)
        
        logger.info(f"文档处理完成: {file.filename}, 生成了 {len(chunks)} 个chunks")
        
        return DocumentProcessResponse(
            success=True,
            message="文档处理成功",
            chunks_count=len(chunks),
            knowledge_id=knowledge_id
        )
        
    except Exception as e:
        logger.error(f"文档处理失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"文档处理失败: {str(e)}")

@app.post("/api/rag/chat", response_model=ChatResponse)
def chat_with_rag(request: ChatRequest):
    """
    RAG智能问答
    从ES检索最相近的top5文档进行回答
    """
    logger.info(f"RAG对话请求: {request.question}, 用户: {request.user_id}")
    
    try:
        # 创建向量存储
        vectorstore = ElasticsearchStore(
            embedding=embeddings,
            index_name=ES_CONFIG['index'],
            es_connection=es_client
        )
        
        # 检索最相近的文档
        docs = vectorstore.similarity_search_with_score(
            request.question,
            k=RAG_CONFIG['top_k']
        )
        
        # 构建知识引用
        references = []
        for doc, score in docs:
            # 计算相关性分数（0-1之间）
            relevance = 1.0 - score  # 距离越小，相关性越高
            
            # 从文档元数据中提取信息
            metadata = doc.metadata
            knowledge_ref = KnowledgeReference(
                knowledge_id=metadata.get("knowledge_id", 0),
                knowledge_name=metadata.get("knowledge_name", "未知知识"),
                description=metadata.get("description", ""),
                tags=metadata.get("tags", []),
                effective_time=metadata.get("effective_time", ""),
                attachments=[metadata.get("source_file", "")],
                relevance=relevance
            )
            references.append(knowledge_ref)
        
        # 生成回答（这里可以集成真实的LLM）
        if references:
            # 基于检索到的文档生成回答
            context = "\n".join([doc.page_content for doc, _ in docs[:3]])  # 使用前3个文档
            answer = f"基于知识库的回答: {request.question}\n\n相关文档内容:\n{context[:RAG_CONFIG['context_limit']]}..."
        else:
            answer = f"抱歉，在知识库中没有找到与'{request.question}'相关的信息。"
        
        return ChatResponse(
            answer=answer,
            references=references
        )
        
    except Exception as e:
        logger.error(f"RAG对话失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"RAG服务异常: {str(e)}")

@app.get("/api/health")
def health_check():
    """健康检查接口"""
    try:
        # 检查ES连接
        es_info = es_client.info()
        return {
            "status": "healthy",
            "elasticsearch": "connected",
            "es_version": es_info.get("version", {}).get("number", "unknown")
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "elasticsearch": "disconnected",
            "error": str(e)
        }

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000) 