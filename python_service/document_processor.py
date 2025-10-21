#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档处理模块
"""

import logging
import os
import tempfile
from pathlib import Path
from typing import List, Dict, Optional, Any
import shutil

# PyMuPDF相关
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("❌ PyMuPDF 不可用")

# Excel解析能力
try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False

# LangChain相关
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False
    print("❌ LangChain 不可用")

# Office文档解析能力（原生）
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    from .utils import (
        convert_with_libreoffice_safe, 
        build_converted_pdf_path, 
        extract_keywords_from_content,
        calculate_chunk_bbox
    )
except ImportError:
    # 当直接运行时使用绝对导入
    import sys
    import os
    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
    from utils import (
        convert_with_libreoffice_safe, 
        build_converted_pdf_path, 
        extract_keywords_from_content,
        calculate_chunk_bbox
    )

logger = logging.getLogger(__name__)




def create_simple_pdf_from_txt(src_path: str, out_path: Optional[str] = None) -> str:
    """TXT 降级方案：将文本简单排版写入单页/多页PDF，便于坐标回显。"""
    with open(src_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    out_pdf = out_path or tempfile.mktemp(suffix='.pdf')
    os.makedirs(os.path.dirname(out_pdf), exist_ok=True)
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    rect = fitz.Rect(40, 50, 555, 800)
    page.insert_textbox(rect, text[:50000], fontsize=11, fontname="helv", align=0)
    doc.save(out_pdf)
    doc.close()
    return out_pdf




def parse_excel_native(file_path: str, filename: str, knowledge_id: int) -> List[Document]:
    """Excel 原生解析，不转PDF"""
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("缺少openpyxl，无法原生解析Excel。请安装: pip install openpyxl")
    wb = load_workbook(file_path, data_only=True)
    chunks: List[Document] = []

    # 构建每个sheet的虚拟坐标网格
    for sheet_idx, ws in enumerate(wb.worksheets):
        # 列宽、行高 → 像素近似（Excel宽度单位转像素，这里做简单近似）
        col_widths = []
        for col in ws.columns:
            # openpyxl列宽在 ws.column_dimensions[col_letter].width
            break
        # 收集列宽（若未设置，给默认 8.43 字符宽 ≈ 64px）
        from openpyxl.utils import get_column_letter
        max_col = ws.max_column
        max_row = ws.max_row
        for c in range(1, max_col + 1):
            letter = get_column_letter(c)
            cw = ws.column_dimensions.get(letter).width if ws.column_dimensions.get(letter) and ws.column_dimensions.get(letter).width else 8.43
            # 近似：1字符宽≈7.5px
            col_widths.append(float(cw) * 7.5)
        row_heights = []
        for r in range(1, max_row + 1):
            rh = ws.row_dimensions.get(r).height if ws.row_dimensions.get(r) and ws.row_dimensions.get(r).height else 15.0
            # 近似：1pt≈1.33px，Excel默认行高≈15pt
            row_heights.append(float(rh) * 1.33)

        # 前缀和得到各列x、各行y起点
        x_starts = [0.0]
        for w in col_widths:
            x_starts.append(x_starts[-1] + w)
        y_starts = [0.0]
        for h in row_heights:
            y_starts.append(y_starts[-1] + h)

        # 构建单元格positions
        positions: List[Dict] = []
        for r in range(1, max_row + 1):
            for c in range(1, max_col + 1):
                cell = ws.cell(row=r, column=c)
                text = str(cell.value) if cell.value is not None else ""
                if not text.strip():
                    continue
                # 合并单元格处理：取所属合并区的外接矩形
                merged_bbox = None
                for rng in ws.merged_cells.ranges:
                    if (r, c) in rng.cells:
                        min_row, min_col, max_row_, max_col_ = rng.min_row, rng.min_col, rng.max_row, rng.max_col
                        x0 = x_starts[min_col - 1]
                        y0 = y_starts[min_row - 1]
                        x1 = x_starts[max_col_]
                        y1 = y_starts[max_row_]
                        merged_bbox = [x0, y0, x1, y1]
                        break
                if merged_bbox:
                    bbox = merged_bbox
                else:
                    x0 = x_starts[c - 1]
                    y0 = y_starts[r - 1]
                    x1 = x_starts[c]
                    y1 = y_starts[r]
                    bbox = [x0, y0, x1, y1]

                positions.append({
                    "text": text.strip(),
                    "bbox": bbox,
                    "sheet": ws.title,
                    "row": r,
                    "col": c,
                })

        # 生成纯文本内容（按行拼接）
        lines = []
        for r in range(1, max_row + 1):
            vals = []
            for c in range(1, max_col + 1):
                val = ws.cell(row=r, column=c).value
                vals.append("" if val is None else str(val))
            lines.append("\t".join(vals).rstrip())
        all_text = "\n".join(lines)

        # 分块（按字符）
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        text_chunks = splitter.split_text(all_text)

        # 将positions按文本包含关系分配给chunk
        for chunk_idx, chunk_text in enumerate(text_chunks):
            chunk_positions = []
            # 简单包含匹配：单元格文本在chunk中出现则纳入（大多数表格有效）
            for p in positions:
                t = p["text"]
                if t and t in chunk_text:
                    chunk_positions.append(p)

            # 计算块级bbox：所有单元格bbox并集（同一sheet）
            if chunk_positions:
                x0 = min(pp["bbox"][0] for pp in chunk_positions)
                y0 = min(pp["bbox"][1] for pp in chunk_positions)
                x1 = max(pp["bbox"][2] for pp in chunk_positions)
                y1 = max(pp["bbox"][3] for pp in chunk_positions)
                bbox = [x0, y0, x1, y1]
            else:
                bbox = [0, 0, 0, 0]

            chunks.append(Document(
                page_content=chunk_text,
                metadata={
                    "knowledge_id": knowledge_id,
                    "source_file": filename,
                    "page_num": sheet_idx + 1,  # 将 sheet 当作"页"
                    "chunk_index": len(chunks),
                    "positions": chunk_positions,
                    "bbox": bbox,
                    "document_name": filename,
                    "document_type": "表格",
                    "sheet_name": ws.title,
                    "keywords": extract_keywords_from_content(chunk_text),
                }
            ))

    return chunks


def parse_txt_native(file_path: str, filename: str, knowledge_id: int) -> List[Document]:
    """TXT 原生解析，不转PDF"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    # 归一化换行
    content = content.replace('\r\n', '\n').replace('\r', '\n')
    lines = content.split('\n')

    # 行级positions（不做真实像素坐标，记录行列与字符区间）
    positions: List[Dict] = []
    global_pos = 0
    for i, line in enumerate(lines, start=1):
        text = line
        start = global_pos
        end = start + len(text)
        positions.append({
            "text": text,
            "line_no": i,
            "char_start": start,
            "char_end": end,
        })
        global_pos = end + 1  # 计入换行

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(content)

    chunks: List[Document] = []
    for idx, chunk_text in enumerate(text_chunks):
        # 为chunk分配行段positions
        chunk_positions = []
        # 通过字符范围粗匹配（行文本在chunk中出现时纳入）
        for p in positions:
            t = p["text"]
            if t and t in chunk_text:
                chunk_positions.append(p)
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": 1,
                "chunk_index": idx,
                "positions": chunk_positions,
                "bbox": [],  # 文本视图通常不需要像素坐标
                "document_name": filename,
                "document_type": "文本",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))

    return chunks


def create_simple_pdf_from_docx(src_path: str, out_path: str) -> tuple[str, List[Dict]]:
    """将DOCX内容简单排版生成PDF，并返回 (pdf_path, paragraph_positions)。
    每个段落在PDF中的放置矩形作为其bbox，供后续可视化使用。
    """
    if not DOCX_AVAILABLE:
        raise RuntimeError("缺少python-docx，无法原生解析Word。请安装: pip install python-docx")

    docx = DocxDocument(src_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    pdf = fitz.open()
    page_width, page_height = 595.0, 842.0  # A4 72dpi points
    margin_left, margin_right, margin_top, margin_bottom = 40.0, 40.0, 50.0, 50.0
    usable_width = page_width - margin_left - margin_right
    cursor_y = margin_top

    paragraph_positions: List[Dict] = []

    # 新建第一页
    page = pdf.new_page(width=page_width, height=page_height)

    def new_page_if_needed(height_needed: float):
        nonlocal page, cursor_y
        if cursor_y + height_needed > (page_height - margin_bottom):
            page = pdf.new_page(width=page_width, height=page_height)
            cursor_y = margin_top

    def draw_paragraph(text: str):
        nonlocal cursor_y
        if not text.strip():
            cursor_y += 8.0
            return
        # 预估高度：用一个较高的矩形，插入后不读取残留，靠换页控制
        rect = fitz.Rect(margin_left, cursor_y, margin_left + usable_width, cursor_y + 200.0)
        new_page_if_needed(rect.height)
        rect = fitz.Rect(margin_left, cursor_y, margin_left + usable_width, cursor_y + 200.0)
        page.insert_textbox(rect, text, fontsize=11, fontname="helv", align=0)
        # 粗略估算占用高度：按文本长度估计行数
        approx_chars_per_line = 90
        lines = max(1, (len(text) // approx_chars_per_line) + 1)
        used_h = 14.0 * lines
        bbox = [rect.x0, rect.y0, rect.x1, rect.y0 + used_h]
        paragraph_positions.append({"text": text, "bbox": bbox, "page": len(pdf)})
        cursor_y += used_h + 6.0

    # 段落
    for p in docx.paragraphs:
        draw_paragraph(p.text)

    # 表格（将每个单元格作为独立小段落）
    for table in docx.tables:
        for row in table.rows:
            cells_text = [c.text.strip() for c in row.cells]
            row_text = "\t".join(cells_text)
            draw_paragraph(row_text)

    pdf.save(out_path)
    pdf.close()
    return out_path, paragraph_positions


def parse_docx_native(file_path: str, filename: str, knowledge_id: int) -> List["Document"]:
    """DOCX 原生解析与PDF生成（不依赖外部工具）"""
    pdf_out = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
    # 生成用于可视化的简单PDF，并获取段落级positions
    _, paragraph_positions = create_simple_pdf_from_docx(file_path, pdf_out)

    # 汇总内容
    all_text = "\n".join([p["text"] for p in paragraph_positions])

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(all_text)

    chunks: List["Document"] = []
    for chunk_idx, chunk_text in enumerate(text_chunks):
        chunk_pos = []
        for p in paragraph_positions:
            t = p.get("text") or ""
            if t and t in chunk_text:
                cp = dict(p)
                # 将页号同步为page_num字段
                cp["page"] = p.get("page", 1)
                chunk_pos.append(cp)
        bbox = calculate_chunk_bbox([{"bbox": pp["bbox"]} for pp in chunk_pos]) if chunk_pos else [0, 0, 0, 0]
        page_counts: Dict[int, int] = {}
        for pp in chunk_pos:
            pg = int(pp.get("page", 1))
            page_counts[pg] = page_counts.get(pg, 0) + 1
        main_page = max(page_counts.items(), key=lambda kv: kv[1])[0] if page_counts else 1
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": main_page,
                "chunk_index": chunk_idx,
                "positions": chunk_pos,
                "bbox": bbox,
                "document_name": filename,
                "document_type": "文档",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))
    return chunks


# ===== PPTX 原生解析与PDF生成（不依赖外部工具） =====
EMU_PER_INCH = 914400.0
POINTS_PER_INCH = 72.0


def create_simple_pdf_from_pptx(src_path: str, out_path: str) -> tuple[str, List[Dict]]:
    """PPTX 原生解析与PDF生成"""
    if not PPTX_AVAILABLE:
        raise RuntimeError("缺少python-pptx，无法原生解析PPT。请安装: pip install python-pptx")

    prs = PptxPresentation(src_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)

    pdf = fitz.open()
    # python-pptx 的尺寸为 Length 类型，用 int() 取 EMU 值更稳妥
    slide_w_pts = (int(prs.slide_width) / EMU_PER_INCH) * POINTS_PER_INCH
    slide_h_pts = (int(prs.slide_height) / EMU_PER_INCH) * POINTS_PER_INCH
    positions: List[Dict] = []

    for s_idx, slide in enumerate(prs.slides):
        page = pdf.new_page(width=slide_w_pts, height=slide_h_pts)
        for shape in slide.shapes:
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = "\n".join([p.text for p in shape.text_frame.paragraphs]).strip()
                if not text:
                    continue
                x0 = (int(shape.left) / EMU_PER_INCH) * POINTS_PER_INCH
                y0 = (int(shape.top) / EMU_PER_INCH) * POINTS_PER_INCH
                x1 = x0 + (int(shape.width) / EMU_PER_INCH) * POINTS_PER_INCH
                y1 = y0 + (int(shape.height) / EMU_PER_INCH) * POINTS_PER_INCH
                rect = fitz.Rect(x0, y0, x1, y1)
                # 在对应矩形内写入文本
                page.insert_textbox(rect, text, fontsize=12, fontname="helv", align=0)
                positions.append({"text": text, "bbox": [x0, y0, x1, y1], "page": s_idx + 1})
            except Exception:
                continue

    pdf.save(out_path)
    pdf.close()
    return out_path, positions


def parse_pptx_native(file_path: str, filename: str, knowledge_id: int) -> List["Document"]:
    """PPTX 原生解析与PDF生成"""
    pdf_out = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
    _, shape_positions = create_simple_pdf_from_pptx(file_path, pdf_out)

    # 汇总内容
    all_text = "\n".join([p["text"] for p in shape_positions])

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
    )
    text_chunks = splitter.split_text(all_text)

    chunks: List["Document"] = []
    for idx, chunk_text in enumerate(text_chunks):
        chunk_pos = []
        for p in shape_positions:
            t = p.get("text") or ""
            if t and t in chunk_text:
                chunk_pos.append(dict(p))
        bbox = calculate_chunk_bbox([{"bbox": pp["bbox"]} for pp in chunk_pos]) if chunk_pos else [0, 0, 0, 0]
        page_counts: Dict[int, int] = {}
        for pp in chunk_pos:
            pg = int(pp.get("page", 1))
            page_counts[pg] = page_counts.get(pg, 0) + 1
        main_page = max(page_counts.items(), key=lambda kv: kv[1])[0] if page_counts else 1
        chunks.append(Document(
            page_content=chunk_text,
            metadata={
                "knowledge_id": knowledge_id,
                "source_file": filename,
                "page_num": main_page,
                "chunk_index": idx,
                "positions": chunk_pos,
                "bbox": bbox,
                "document_name": filename,
                "document_type": "演示文稿",
                "keywords": extract_keywords_from_content(chunk_text),
            }
        ))
    return chunks


def process_document_unified(
    file_path: str,
    filename: str,
    knowledge_id: int,
    knowledge_name: Optional[str] = None,
    description: Optional[str] = None,
    tags: Optional[str] = None,
    effective_time: Optional[str] = None,
    workspaces: Optional[str] = None,
) -> Dict:
    """
    统一处理文档：
    - 非PDF先转换为PDF（优先使用LibreOffice），不改动原始文件
    - 基于PDF用PyMuPDF提取块级坐标并分块
    """
    try:
        # 1) 非PDF → PDF（不改动原始文件，仅生成临时PDF用于定位与回显）
        ext = Path(filename).suffix.lower()
        pdf_path_to_open = file_path
        temp_generated_pdf = None

        if ext in {'.xlsx', '.xls'}:
            # Excel 原生解析
            chunks = parse_excel_native(file_path, filename, knowledge_id)
        elif ext == '.txt':
            # TXT 原生解析
            chunks = parse_txt_native(file_path, filename, knowledge_id)
        elif ext in {'.doc', '.docx'}:
            # Word 原生解析
            chunks = parse_docx_native(file_path, filename, knowledge_id)
            # 使用我们生成的简易PDF供回显
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            pdf_path_to_open = target_pdf
        elif ext in {'.ppt', '.pptx'}:
            # PPT 原生解析
            chunks = parse_pptx_native(file_path, filename, knowledge_id)
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            pdf_path_to_open = target_pdf
        elif ext != '.pdf':
            # 其他Office类型仍尝试转为PDF（如需）
            target_pdf = build_converted_pdf_path(int(knowledge_id) if knowledge_id else 0, filename)
            try:
                target_dir = os.path.dirname(target_pdf)
                src_for_convert = file_path
                pdf_path_to_open = convert_with_libreoffice_safe(src_for_convert, out_dir=target_dir)
                if os.path.abspath(pdf_path_to_open) != os.path.abspath(target_pdf):
                    shutil.copyfile(pdf_path_to_open, target_pdf)
                    pdf_path_to_open = target_pdf
                logger.info(f"已将 {filename} 转为PDF: {pdf_path_to_open}")
            except Exception as e:
                if ext == '.txt':
                    pdf_path_to_open = create_simple_pdf_from_txt(file_path, out_path=target_pdf)
                    logger.info(f"TXT降级转PDF成功: {pdf_path_to_open}")
                else:
                    raise

        # 2) 基于PDF走统一解析 - 使用新的匹配切分策略
        if ext == '.pdf' or pdf_path_to_open:
            doc = fitz.open(pdf_path_to_open)
            try:
                logger.info(f"成功打开文档，页数: {len(doc)}")

                # 使用新的PDF切分策略，确保大块和小块完全匹配
                try:
                    from .pdf_chunking_strategy import process_pdf_with_matched_chunks
                except ImportError:
                    from pdf_chunking_strategy import process_pdf_with_matched_chunks
                
                # 使用更小的块大小，确保不超过token限制
                documents = process_pdf_with_matched_chunks(doc, filename, max_chunk_size=1500, overlap=100)
                logger.info(f"PDF切分后生成 {len(documents)} 个chunks")

                chunks = []
                for doc in documents:
                    # 为每个chunk添加知识库元数据
                    doc.metadata.update({
                        "knowledge_id": knowledge_id,
                        "knowledge_name": knowledge_name or "",
                        "description": description or "",
                        "tags": tags or "",
                        "effective_time": effective_time or "",
                        "document_name": filename,
                        "keywords": extract_keywords_from_content(doc.page_content),
                    })
                    chunks.append(doc)
            finally:
                try:
                    doc.close()
                except Exception:
                    pass
                # 不删除持久化PDF，供前端下载/回显
        
        # 为Excel/TXT等原生解析生成的chunk补齐知识元数据字段
        try:
            for ch in chunks:
                md = ch.metadata
                if "knowledge_name" not in md:
                    md["knowledge_name"] = knowledge_name or ""
                if "description" not in md:
                    md["description"] = description or ""
                if "tags" not in md:
                    md["tags"] = tags or ""
                if "effective_time" not in md:
                    md["effective_time"] = effective_time or ""
        except Exception:
            pass

        logger.info(f"最终生成 {len(chunks)} 个chunks")
        
        # 存储到ES
        workspaces_list = workspaces.split(",") if workspaces else None
        try:
            from .es_client import store_chunks_to_es
        except ImportError:
            from es_client import store_chunks_to_es
        stored_count = store_chunks_to_es(chunks, knowledge_id, workspaces_list)
        
        return {
            "chunks_count": stored_count,
            "total_chunks": len(chunks),
            "success": stored_count > 0
        }
        
    except Exception as e:
        logger.error(f"文档处理失败: {e}")
        import traceback
        traceback.print_exc()
        raise e
