#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件切块测试脚本
测试不同文件类型的切块效果和性能
"""

import os
import sys
import time
import json
from pathlib import Path
from typing import List, Dict, Any
import tempfile

# 添加python_service目录到路径
sys.path.append('python_service')

# 导入文档处理相关库
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# 导入文档处理库
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# 导入配置
from python_service.config import DOCUMENT_CONFIG, CHUNKING_CONFIG

class ChunkingTester:
    def __init__(self):
        self.results = {}
        self.test_files = {
            'pdf': ['python_service/file/安联美元.pdf', 'python_service/file/店铺入住流程.pdf'],
            'docx': ['python_service/file/manus介绍.docx'],
            'xlsx': ['python_service/file/小红书选品.xlsx'],
            'pptx': ['python_service/file/network_skill.pptx']
        }
        
    def test_pdf_chunking(self, file_path: str) -> Dict[str, Any]:
        """测试PDF文件切块"""
        print(f"测试PDF文件: {file_path}")
        
        start_time = time.time()
        
        try:
            # 使用PyMuPDFLoader处理PDF
            loader = PyMuPDFLoader(file_path)
            documents = loader.load()
            
            # 文本分割
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_documents(documents)
            
            end_time = time.time()
            
            # 统计信息
            total_text_length = sum(len(chunk.page_content) for chunk in chunks)
            avg_chunk_length = total_text_length / len(chunks) if chunks else 0
            
            # 分析chunk内容
            chunk_analysis = []
            for i, chunk in enumerate(chunks[:5]):  # 只分析前5个chunk
                chunk_analysis.append({
                    'chunk_index': i,
                    'length': len(chunk.page_content),
                    'preview': chunk.page_content[:100] + '...' if len(chunk.page_content) > 100 else chunk.page_content,
                    'metadata': chunk.metadata
                })
            
            return {
                'success': True,
                'processing_time': end_time - start_time,
                'total_chunks': len(chunks),
                'total_text_length': total_text_length,
                'avg_chunk_length': avg_chunk_length,
                'chunk_analysis': chunk_analysis,
                'file_size_mb': os.path.getsize(file_path) / (1024 * 1024)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'processing_time': time.time() - start_time
            }
    
    def test_docx_chunking(self, file_path: str) -> Dict[str, Any]:
        """测试DOCX文件切块"""
        print(f"测试DOCX文件: {file_path}")
        
        start_time = time.time()
        
        try:
            # 处理Word文档
            doc = DocxDocument(file_path)
            text_parts = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
            end_time = time.time()
            
            # 统计信息
            total_text_length = sum(len(chunk.page_content) for chunk in chunks)
            avg_chunk_length = total_text_length / len(chunks) if chunks else 0
            
            # 分析chunk内容
            chunk_analysis = []
            for i, chunk in enumerate(chunks[:5]):  # 只分析前5个chunk
                chunk_analysis.append({
                    'chunk_index': i,
                    'length': len(chunk.page_content),
                    'preview': chunk.page_content[:100] + '...' if len(chunk.page_content) > 100 else chunk.page_content
                })
            
            return {
                'success': True,
                'processing_time': end_time - start_time,
                'total_chunks': len(chunks),
                'total_text_length': total_text_length,
                'avg_chunk_length': avg_chunk_length,
                'chunk_analysis': chunk_analysis,
                'file_size_mb': os.path.getsize(file_path) / (1024 * 1024)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'processing_time': time.time() - start_time
            }
    
    def test_xlsx_chunking(self, file_path: str) -> Dict[str, Any]:
        """测试XLSX文件切块"""
        print(f"测试XLSX文件: {file_path}")
        
        start_time = time.time()
        
        try:
            # 处理Excel文件
            workbook = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"工作表: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None and str(cell_value).strip():
                            row_text.append(str(cell_value).strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
            end_time = time.time()
            
            # 统计信息
            total_text_length = sum(len(chunk.page_content) for chunk in chunks)
            avg_chunk_length = total_text_length / len(chunks) if chunks else 0
            
            # 分析chunk内容
            chunk_analysis = []
            for i, chunk in enumerate(chunks[:5]):  # 只分析前5个chunk
                chunk_analysis.append({
                    'chunk_index': i,
                    'length': len(chunk.page_content),
                    'preview': chunk.page_content[:100] + '...' if len(chunk.page_content) > 100 else chunk.page_content
                })
            
            return {
                'success': True,
                'processing_time': end_time - start_time,
                'total_chunks': len(chunks),
                'total_text_length': total_text_length,
                'avg_chunk_length': avg_chunk_length,
                'chunk_analysis': chunk_analysis,
                'file_size_mb': os.path.getsize(file_path) / (1024 * 1024)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'processing_time': time.time() - start_time
            }
    
    def test_pptx_chunking(self, file_path: str) -> Dict[str, Any]:
        """测试PPTX文件切块"""
        print(f"测试PPTX文件: {file_path}")
        
        start_time = time.time()
        
        try:
            # 处理PowerPoint文件
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_number, slide in enumerate(prs.slides, 1):
                text_parts.append(f"幻灯片 {slide_number}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                text_parts.append("")  # 空行分隔幻灯片
            
            text = "\n".join(text_parts)
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=DOCUMENT_CONFIG['chunk_size'],
                chunk_overlap=DOCUMENT_CONFIG['chunk_overlap'],
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            chunks = [Document(page_content=chunk) for chunk in chunks]
            
            end_time = time.time()
            
            # 统计信息
            total_text_length = sum(len(chunk.page_content) for chunk in chunks)
            avg_chunk_length = total_text_length / len(chunks) if chunks else 0
            
            # 分析chunk内容
            chunk_analysis = []
            for i, chunk in enumerate(chunks[:5]):  # 只分析前5个chunk
                chunk_analysis.append({
                    'chunk_index': i,
                    'length': len(chunk.page_content),
                    'preview': chunk.page_content[:100] + '...' if len(chunk.page_content) > 100 else chunk.page_content
                })
            
            return {
                'success': True,
                'processing_time': end_time - start_time,
                'total_chunks': len(chunks),
                'total_text_length': total_text_length,
                'avg_chunk_length': avg_chunk_length,
                'chunk_analysis': chunk_analysis,
                'file_size_mb': os.path.getsize(file_path) / (1024 * 1024)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'processing_time': time.time() - start_time
            }
    
    def run_all_tests(self):
        """运行所有测试"""
        print("=" * 60)
        print("开始文件切块测试")
        print("=" * 60)
        
        # 测试PDF文件
        for pdf_file in self.test_files['pdf']:
            if os.path.exists(pdf_file):
                result = self.test_pdf_chunking(pdf_file)
                self.results[f"PDF_{Path(pdf_file).name}"] = result
                self.print_result("PDF", Path(pdf_file).name, result)
        
        # 测试DOCX文件
        for docx_file in self.test_files['docx']:
            if os.path.exists(docx_file):
                result = self.test_docx_chunking(docx_file)
                self.results[f"DOCX_{Path(docx_file).name}"] = result
                self.print_result("DOCX", Path(docx_file).name, result)
        
        # 测试XLSX文件
        for xlsx_file in self.test_files['xlsx']:
            if os.path.exists(xlsx_file):
                result = self.test_xlsx_chunking(xlsx_file)
                self.results[f"XLSX_{Path(xlsx_file).name}"] = result
                self.print_result("XLSX", Path(xlsx_file).name, result)
        
        # 测试PPTX文件
        for pptx_file in self.test_files['pptx']:
            if os.path.exists(pptx_file):
                result = self.test_pptx_chunking(pptx_file)
                self.results[f"PPTX_{Path(pptx_file).name}"] = result
                self.print_result("PPTX", Path(pptx_file).name, result)
        
        # 生成测试报告
        self.generate_report()
    
    def print_result(self, file_type: str, filename: str, result: Dict[str, Any]):
        """打印测试结果"""
        print(f"\n{'-' * 40}")
        print(f"文件类型: {file_type}")
        print(f"文件名: {filename}")
        
        if result['success']:
            print(f"✅ 处理成功")
            print(f"📊 处理时间: {result['processing_time']:.2f}秒")
            print(f"📄 文件大小: {result['file_size_mb']:.2f}MB")
            print(f"🔢 总块数: {result['total_chunks']}")
            print(f"📏 总文本长度: {result['total_text_length']}字符")
            print(f"📐 平均块长度: {result['avg_chunk_length']:.1f}字符")
            
            if result['chunk_analysis']:
                print(f"\n📋 前5个块分析:")
                for chunk_info in result['chunk_analysis']:
                    print(f"  块{chunk_info['chunk_index']}: {chunk_info['length']}字符")
                    print(f"    预览: {chunk_info['preview']}")
        else:
            print(f"❌ 处理失败")
            print(f"🚨 错误信息: {result['error']}")
    
    def generate_report(self):
        """生成测试报告"""
        print(f"\n{'=' * 60}")
        print("测试报告总结")
        print(f"{'=' * 60}")
        
        # 统计成功和失败
        success_count = sum(1 for result in self.results.values() if result['success'])
        total_count = len(self.results)
        
        print(f"📈 总测试文件数: {total_count}")
        print(f"✅ 成功处理: {success_count}")
        print(f"❌ 失败处理: {total_count - success_count}")
        
        if success_count > 0:
            # 计算平均处理时间
            avg_time = sum(result['processing_time'] for result in self.results.values() if result['success']) / success_count
            print(f"⏱️ 平均处理时间: {avg_time:.2f}秒")
            
            # 按文件类型统计
            type_stats = {}
            for key, result in self.results.items():
                if result['success']:
                    file_type = key.split('_')[0]
                    if file_type not in type_stats:
                        type_stats[file_type] = []
                    type_stats[file_type].append(result)
            
            print(f"\n📊 按文件类型统计:")
            for file_type, results in type_stats.items():
                avg_chunks = sum(r['total_chunks'] for r in results) / len(results)
                avg_length = sum(r['avg_chunk_length'] for r in results) / len(results)
                print(f"  {file_type}: 平均{avg_chunks:.1f}块, 平均块长度{avg_length:.1f}字符")
        
        # 保存详细报告
        report_file = "chunking_test_report.json"
        with open(report_file, 'w', encoding='utf-8') as f:
            json.dump(self.results, f, ensure_ascii=False, indent=2)
        print(f"\n📄 详细报告已保存到: {report_file}")

def main():
    """主函数"""
    tester = ChunkingTester()
    tester.run_all_tests()

if __name__ == "__main__":
    main() 