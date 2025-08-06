#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档处理功能测试脚本
测试Python服务的文档处理能力
"""

import requests
import json
import os
from pathlib import Path

# Python服务配置
PYTHON_SERVICE_URL = "http://localhost:8000"

def test_python_health():
    """测试Python服务健康状态"""
    try:
        response = requests.get(f"{PYTHON_SERVICE_URL}/api/health")
        print(f"Python服务健康检查: {response.status_code}")
        if response.status_code == 200:
            print("✅ Python服务正常运行")
            return True
        else:
            print("❌ Python服务异常")
            return False
    except Exception as e:
        print(f"❌ 无法连接到Python服务: {e}")
        return False

def test_document_processing(file_path, knowledge_id=1):
    """测试文档处理功能"""
    if not os.path.exists(file_path):
        print(f"❌ 文件不存在: {file_path}")
        return False
    
    file_name = os.path.basename(file_path)
    file_extension = Path(file_path).suffix.lower()
    
    print(f"\n📄 测试文档处理: {file_name}")
    print(f"文件类型: {file_extension}")
    
    try:
        # 准备multipart数据
        with open(file_path, 'rb') as f:
            files = {'file': (file_name, f, 'application/octet-stream')}
            data = {
                'knowledge_id': knowledge_id,
                'knowledge_name': f'测试知识_{file_name}',
                'description': f'测试文档处理功能 - {file_name}',
                'tags': '测试,文档处理',
                'effective_time': '2024-01-01 至 2024-12-31'
            }
            
            # 发送请求
            response = requests.post(
                f"{PYTHON_SERVICE_URL}/api/document/process",
                files=files,
                data=data,
                timeout=60
            )
            
            print(f"响应状态码: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                print("✅ 文档处理成功")
                print(f"生成的chunks数量: {result.get('chunks_count', 0)}")
                print(f"知识ID: {result.get('knowledge_id')}")
                return True
            else:
                print(f"❌ 文档处理失败: {response.text}")
                return False
                
    except Exception as e:
        print(f"❌ 文档处理异常: {e}")
        return False

def test_rag_chat():
    """测试RAG对话功能"""
    print("\n🤖 测试RAG对话功能")
    
    try:
        data = {
            "question": "测试问题：请介绍一下文档处理功能",
            "user_id": "test_user"
        }
        
        response = requests.post(
            f"{PYTHON_SERVICE_URL}/api/rag/chat",
            json=data,
            timeout=30
        )
        
        print(f"响应状态码: {response.status_code}")
        
        if response.status_code == 200:
            result = response.json()
            print("✅ RAG对话成功")
            print(f"答案: {result.get('answer', '')}")
            print(f"引用数量: {len(result.get('references', []))}")
            return True
        else:
            print(f"❌ RAG对话失败: {response.text}")
            return False
            
    except Exception as e:
        print(f"❌ RAG对话异常: {e}")
        return False

def create_test_files():
    """创建测试文件"""
    test_files = []
    
    # 创建测试Word文档
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('测试Word文档', 0)
        doc.add_paragraph('这是一个测试Word文档，用于验证文档处理功能。')
        doc.add_paragraph('文档包含多个段落和表格。')
        
        # 添加表格
        table = doc.add_table(rows=3, cols=3)
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                cell.text = f'单元格{i+1}-{j+1}'
        
        word_file = 'test_document.docx'
        doc.save(word_file)
        test_files.append(word_file)
        print(f"✅ 创建测试Word文档: {word_file}")
    except Exception as e:
        print(f"❌ 创建Word文档失败: {e}")
    
    # 创建测试Excel文件
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "测试工作表"
        
        # 添加数据
        ws['A1'] = '姓名'
        ws['B1'] = '年龄'
        ws['C1'] = '部门'
        ws['A2'] = '张三'
        ws['B2'] = 25
        ws['C2'] = '技术部'
        ws['A3'] = '李四'
        ws['B3'] = 30
        ws['C3'] = '产品部'
        
        excel_file = 'test_document.xlsx'
        wb.save(excel_file)
        test_files.append(excel_file)
        print(f"✅ 创建测试Excel文件: {excel_file}")
    except Exception as e:
        print(f"❌ 创建Excel文件失败: {e}")
    
    # 创建测试PowerPoint文件
    try:
        from pptx import Presentation
        prs = Presentation()
        
        # 添加标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "测试演示文稿"
        subtitle.text = "用于验证文档处理功能"
        
        # 添加内容页
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "功能特点"
        tf = body_shape.text_frame
        tf.text = "支持多种文档格式"
        p = tf.add_paragraph()
        p.text = "智能文本提取"
        p = tf.add_paragraph()
        p.text = "向量化存储"
        
        pptx_file = 'test_document.pptx'
        prs.save(pptx_file)
        test_files.append(pptx_file)
        print(f"✅ 创建测试PowerPoint文件: {pptx_file}")
    except Exception as e:
        print(f"❌ 创建PowerPoint文件失败: {e}")
    
    # 创建测试文本文件
    try:
        txt_file = 'test_document.txt'
        with open(txt_file, 'w', encoding='utf-8') as f:
            f.write("这是一个测试文本文件。\n")
            f.write("用于验证文档处理功能。\n")
            f.write("支持UTF-8编码的中文内容。\n")
            f.write("包含多行文本内容。\n")
        
        test_files.append(txt_file)
        print(f"✅ 创建测试文本文件: {txt_file}")
    except Exception as e:
        print(f"❌ 创建文本文件失败: {e}")
    
    return test_files

def cleanup_test_files(test_files):
    """清理测试文件"""
    print("\n🧹 清理测试文件...")
    for file_path in test_files:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"✅ 删除测试文件: {file_path}")
        except Exception as e:
            print(f"❌ 删除文件失败 {file_path}: {e}")

def main():
    """主测试函数"""
    print("🚀 开始文档处理功能测试")
    print("=" * 50)
    
    # 1. 测试Python服务健康状态
    if not test_python_health():
        print("❌ Python服务不可用，测试终止")
        return
    
    # 2. 创建测试文件
    print("\n📝 创建测试文件...")
    test_files = create_test_files()
    
    if not test_files:
        print("❌ 无法创建测试文件，测试终止")
        return
    
    # 3. 测试文档处理
    print("\n📄 测试文档处理功能...")
    success_count = 0
    total_count = len(test_files)
    
    for file_path in test_files:
        if test_document_processing(file_path):
            success_count += 1
    
    # 4. 测试RAG对话
    test_rag_chat()
    
    # 5. 清理测试文件
    cleanup_test_files(test_files)
    
    # 6. 输出测试结果
    print("\n" + "=" * 50)
    print("📊 测试结果总结")
    print(f"文档处理成功率: {success_count}/{total_count}")
    
    if success_count == total_count:
        print("🎉 所有测试通过！")
    else:
        print("⚠️  部分测试失败，请检查服务配置")

if __name__ == "__main__":
    main() 